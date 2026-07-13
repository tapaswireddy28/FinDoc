"""Document parsing: extract clean, page-aware text from many file types.

Supported: PDF, TXT/MD, PPTX (one "page" per slide), and images (PNG/JPG/...
transcribed with Claude vision via the optional `ocr` callback). The "page"
field is what gets cited back with answers.
"""
import re
from pathlib import Path

from config import TEXT_EXTS, PDF_EXTS, PPTX_EXTS, IMAGE_EXTS, SUPPORTED_EXTS


def clean_text(text: str) -> str:
    """Collapse whitespace and strip control characters."""
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)        # collapse runs of spaces/tabs
    text = re.sub(r"\n{3,}", "\n\n", text)     # collapse big blank gaps
    return text.strip()


def extract_pages(path: str | Path, ocr=None) -> list[dict]:
    """Return a list of {"page": int, "text": str}, one entry per page/slide.

    `ocr` is an optional callable(path) -> str used to transcribe image files
    (e.g. Generator.ocr_image). Without it, images cannot be read.
    """
    path = Path(path)
    if not path.exists():
        raise FileNotFoundError(f"Document not found: {path}")

    suffix = path.suffix.lower()

    if suffix in TEXT_EXTS:
        text = clean_text(path.read_text(encoding="utf-8", errors="ignore"))
        return [{"page": 1, "text": text}] if text else []

    if suffix in PDF_EXTS:
        import fitz  # PyMuPDF (imported here so non-PDF use needs no PDF lib)

        pages = []
        with fitz.open(path) as doc:
            for i, page in enumerate(doc, start=1):
                text = clean_text(page.get_text("text"))
                if text:
                    pages.append({"page": i, "text": text})
        return pages

    if suffix in PPTX_EXTS:
        from pptx import Presentation  # python-pptx

        pages = []
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, start=1):
            parts = []
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    parts.append(shape.text_frame.text)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text for c in row.cells))
            text = clean_text("\n".join(parts))
            if text:
                pages.append({"page": i, "text": text})  # page == slide number
        return pages

    if suffix in IMAGE_EXTS:
        if ocr is None:
            raise ValueError(
                "Image files need Claude vision OCR, which requires "
                "LLMFOUNDRY_TOKEN to be set."
            )
        text = clean_text(ocr(path) or "")
        return [{"page": 1, "text": text}] if text else []

    raise ValueError(
        f"Unsupported file type '{suffix}'. Supported: "
        + ", ".join(sorted(SUPPORTED_EXTS))
    )
